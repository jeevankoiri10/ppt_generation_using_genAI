import io
import os

from django.core.files.base import ContentFile

from dash.models import GenerationHistory

load_ai = os.getenv("LOADAI", "False") == "True"

if load_ai:
    import torch
    import pandas as pd
    from pptx import Presentation
    from keras.preprocessing.text import Tokenizer
    from keras.utils import pad_sequences
    from transformers import AutoModelForSeq2SeqLM, AutoTokenizer

    device = "cuda" if torch.cuda.is_available() else "cpu"
    model_ckpt = "./assets/summarization_model/cnn_dailymail"
    tokenizer = AutoTokenizer.from_pretrained(model_ckpt)
    model_pegasus = AutoModelForSeq2SeqLM.from_pretrained(model_ckpt)

    from transformers import pipeline

    pipe = pipeline('summarization', model=model_pegasus, tokenizer=tokenizer)

    data = pd.read_csv('./assets/sentences.txt', sep=',', header=None)
    text_data = data[1]

    from keras.models import load_model

    max_words = 10000
    tokenizer = Tokenizer(num_words=max_words)
    tokenizer.fit_on_texts(text_data)
    loaded_model = load_model('./assets/heading_classifier_model.h5')


def classify_heading(sentence):
    sentence_sequence = tokenizer.texts_to_sequences([sentence])
    sentence_sequence = pad_sequences(sentence_sequence, maxlen=15)
    prediction = loaded_model.predict(sentence_sequence)

    if prediction > 0.5:
        return "Heading"
    else:
        return "Not Heading"


def slideGenerator(title, subtitle, presentation):
    slide_layout = presentation.slide_layouts[1]
    """ Ref for different slide layouts on the basis of index: 
    0 -> title and subtitle => Used For the main Topic and the author name 
    1 -> title and content => Used For the sub topic with the summary
    """
    slide = presentation.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def process_with_ai(generation_request: GenerationHistory):
    if not load_ai:
        return

    input_file = generation_request.document
    with open(input_file.path, 'r') as file:
        presentation = Presentation("./assets/template3.pptx")
        content = file.read()
        title = generation_request.title
        author = generation_request.author_override if generation_request.author_override else generation_request.author.username

        lines = content.split('\n')

        lines_array = []

        for line in lines:
            lines_array.append(line)

        lines_array = [line.strip() for line in lines_array if line.strip()]

        json_objects = []
        headings_array = []
        content = ''
        for line in lines_array:
            result = classify_heading(line)
            if result == 'Heading':
                if content:
                    summary = pipe(content)[0]['summary_text']
                    json_objects.append({'heading': headings_array[-1], 'summary': summary})
                # Start a new content string
                content = ''
                headings_array.append(line)
            else:
                content += line + ' '
        if content:
            summary = pipe(content)[0]['summary_text']
            json_objects.append({'heading': headings_array[-1], 'summary': summary})

        slide_layout0 = presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(slide_layout0)
        slide.shapes.title.text = 'A Presentation On\n' + title
        slide.placeholders[1].text = author

        for element in json_objects:
            text = element['summary']
            summaryconverted = text.replace(". ", ".\n")

            slideGenerator(element['heading'], summaryconverted, presentation)

        output = io.BytesIO()
        presentation.save(output)

        content_file = ContentFile(output.getvalue())

        generation_request.presentation.save(f'{generation_request.pk}_output_generation.pptx', content_file)

        generation_request.save()
